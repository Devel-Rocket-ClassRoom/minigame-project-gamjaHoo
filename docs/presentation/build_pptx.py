# -*- coding: utf-8 -*-
"""
FM-Lite 발표 자료 (PPTX) 생성 스크립트.
- 다크(게임 무드) 테마 / 한국어 / 16:9
- 스크린샷은 플레이스홀더 박스 + 별도 캡처 가이드(capture-guide.md)
실행: python docs/presentation/build_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")
FRAME = RGBColor(0x3A, 0x44, 0x63)

# ---- 팔레트 (다크) ----
BG      = RGBColor(0x12, 0x14, 0x1F)   # 메인 배경 (짙은 네이비)
BG_CARD = RGBColor(0x1C, 0x22, 0x38)   # 카드/패널
BG_CODE = RGBColor(0x0D, 0x10, 0x1A)   # 코드 블록
ACCENT  = RGBColor(0x4A, 0x90, 0xD9)   # 게임 액센트 블루
GREEN   = RGBColor(0x2E, 0xCC, 0x71)   # 개선/긍정
AMBER   = RGBColor(0xF3, 0x9C, 0x12)   # 주의/난관
TEXT    = RGBColor(0xEC, 0xEC, 0xF1)   # 본문
MUTED   = RGBColor(0x9A, 0xA3, 0xB6)   # 보조
LINE    = RGBColor(0x2A, 0x31, 0x4C)

FONT = "맑은 고딕"
MONO = "Consolas"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


# ---------- 헬퍼 ----------
def _set_font(run, size=18, bold=False, color=TEXT, name=FONT, italic=False):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = name
    # 한글 글리프용 East-Asian / Complex-Script typeface 명시
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def rect(s, l, t, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, l, t, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def textbox(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    return tf


def para(tf, runs, align=PP_ALIGN.LEFT, space_after=8, space_before=0, line=1.05, first=False, bullet=None, level=0):
    """runs = [(text, dict), ...]"""
    p = tf.paragraphs[0] if first and tf.paragraphs[0].text == "" and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    try:
        p.line_spacing = line
    except Exception:
        pass
    p.level = level
    if bullet is not None:
        r = p.add_run(); r.text = bullet + "  "
        _set_font(r, **{"size": runs[0][1].get("size", 18), "bold": True, "color": runs[0][1].get("bcolor", ACCENT)})
    for text, kw in runs:
        kw = dict(kw); kw.pop("bcolor", None)
        r = p.add_run(); r.text = text
        _set_font(r, **kw)
    return p


def header(s, kicker, title):
    """상단 키커 + 제목 + 액센트 바"""
    rect(s, Inches(0.7), Inches(0.62), Inches(0.14), Inches(0.62), fill=ACCENT)
    tf = textbox(s, Inches(0.95), Inches(0.5), Inches(11.6), Inches(1.1))
    para(tf, [(kicker, {"size": 13, "bold": True, "color": ACCENT})], space_after=2, first=True)
    para(tf, [(title, {"size": 30, "bold": True, "color": TEXT})], space_after=0)
    rect(s, Inches(0.7), Inches(1.62), Inches(11.93), Pt(1.4), fill=LINE)


def footer(s, n):
    tf = textbox(s, Inches(11.4), Inches(7.02), Inches(1.6), Inches(0.35))
    para(tf, [("FM-Lite", {"size": 9, "color": MUTED}), (f"   ·   {n:02d}", {"size": 9, "color": ACCENT, "bold": True})],
         align=PP_ALIGN.RIGHT, first=True)


def placeholder(s, l, t, w, h, caption):
    box = rect(s, l, t, w, h, fill=BG_CARD, line=ACCENT, line_w=1.25, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    box.line.dash_style  # noop guard
    ln = box.line._get_or_add_ln()
    d = ln.makeelement(qn("a:prstDash"), {"val": "dash"}); ln.append(d)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, [("📷  스크린샷", {"size": 14, "bold": True, "color": ACCENT})], align=PP_ALIGN.CENTER, first=True, space_after=4)
    para(tf, [(caption, {"size": 11, "color": MUTED})], align=PP_ALIGN.CENTER, space_after=0)


def card(s, l, t, w, h, fill=BG_CARD, accent=None):
    c = rect(s, l, t, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    if accent:
        rect(s, l, t, Inches(0.09), h, fill=accent, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    return c


def img_path(name):
    p = os.path.join(ASSETS, name)
    return p if os.path.exists(p) else None


def image(s, name, l, t, w, caption=None):
    """assets/name 이미지를 폭 w(inch) 로 비율 유지 배치. 파일 없으면 None 반환(폴백용)."""
    from PIL import Image as _PILImage
    p = img_path(name)
    if p is None:
        return None
    iw, ih = _PILImage.open(p).size
    h = w * ih / iw
    rect(s, Inches(l-0.035), Inches(t-0.035), Inches(w+0.07), Inches(h+0.07),
         fill=None, line=FRAME, line_w=1.25)
    s.shapes.add_picture(p, Inches(l), Inches(t), Inches(w), Inches(h))
    if caption:
        tf = textbox(s, Inches(l), Inches(t+h+0.07), Inches(w), Inches(0.3))
        para(tf, [(caption, {"size": 10.5, "italic": True, "color": MUTED})], align=PP_ALIGN.CENTER, first=True)
    return h


def image_or_placeholder(s, name, l, t, w, h_box, caption):
    """이미지가 있으면 박스 중앙에 배치, 없으면 점선 플레이스홀더."""
    from PIL import Image as _PILImage
    p = img_path(name)
    if p is None:
        placeholder(s, Inches(l), Inches(t), Inches(w), Inches(h_box), caption)
        return
    iw, ih = _PILImage.open(p).size
    h = w * ih / iw
    ty = t + max(0.0, (h_box - h - 0.32) / 2)
    image(s, name, l, ty, w, caption)


N = [0]
def pageno():
    N[0] += 1
    return N[0]


# =========================================================
# 1. 타이틀
# =========================================================
s = slide()
# 배경 액센트 도형
rect(s, Inches(9.5), Inches(-1.7), Inches(5.6), Inches(5.6), fill=BG_CARD, shape=MSO_SHAPE.OVAL)
rect(s, Inches(1.1), Inches(2.0), Inches(0.16), Inches(3.0), fill=ACCENT)
image(s, "01_main_menu.png", 7.45, 2.45, 5.4, "메인 메뉴")
tf = textbox(s, Inches(1.5), Inches(2.05), Inches(5.7), Inches(3.2))
para(tf, [("UNITY · C# · 3주 1인 개발", {"size": 15, "bold": True, "color": ACCENT})], space_after=10, first=True)
para(tf, [("FM-Lite", {"size": 66, "bold": True, "color": TEXT})], space_after=6)
para(tf, [("FM의 ‘육성’만 떼어낸, 한 시즌 2~3시간 컴팩트 축구 매니저", {"size": 20, "color": MUTED})], space_after=0)
tf2 = textbox(s, Inches(1.5), Inches(6.3), Inches(10), Inches(0.6))
para(tf2, [("발표자: JihooKim", {"size": 13, "color": TEXT}),
           ("      2026.06.07      프로젝트 마무리 발표", {"size": 13, "color": MUTED})], first=True)

# =========================================================
# 2. 목차
# =========================================================
s = slide(); header(s, "AGENDA", "발표 순서")
items = [
    ("01", "게임 소개 & 기술 스택", "무엇을, 무엇으로 만들었나"),
    ("02", "구현 범위", "완성된 핵심 시스템 맵"),
    ("03", "개발 과정", "설계가 바뀐 순간들 · 난관과 극복"),
    ("04", "느낀 점", "3주가 남긴 것"),
    ("05", "대표 구현", "매치 시뮬레이션 · 정보 비대칭 · AI 협업"),
    ("06", "Q & A", ""),
]
y = 1.95
for num, title, desc in items:
    card(s, Inches(0.95), Inches(y), Inches(11.4), Inches(0.72), accent=ACCENT)
    tf = textbox(s, Inches(1.25), Inches(y), Inches(11), Inches(0.72), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(num, {"size": 18, "bold": True, "color": ACCENT}),
              ("    " + title, {"size": 18, "bold": True, "color": TEXT}),
              (("      —  " + desc) if desc else "", {"size": 13, "color": MUTED})], first=True, space_after=0)
    y += 0.83
footer(s, pageno())

# =========================================================
# 3. 기술 스택
# =========================================================
s = slide(); header(s, "TECH STACK", "기술 스택")
cols = [
    ("엔진 · 언어", [("Unity 6000.3", "C# / IL2CPP 빌드"),
                  ("ScriptableObject", "정적 데이터 · 밸런스 외부화"),
                  ("Newtonsoft.Json", "세이브 슬롯 직렬화")]),
    ("UI · 텍스트", [("Modern UI Pack", "버튼·드롭다운·모달·툴팁"),
                 ("TextMeshPro", "리치 텍스트 · 다국어 입력"),
                 ("NotoSansKR (OFL)", "한/영 가변 폰트")]),
    ("구조 · 협업", [("4-Layer + EventBus", "Domain/Core/App/Persistence"),
                  ("Claude Code + Unity MCP", "코드·씬·PR 자동화"),
                  ("GitHub Issues / Projects", "이슈→브랜치→PR→머지")]),
]
x = 0.95
for ctitle, rows in cols:
    card(s, Inches(x), Inches(2.0), Inches(3.66), Inches(4.1))
    tf = textbox(s, Inches(x+0.28), Inches(2.25), Inches(3.1), Inches(3.7))
    para(tf, [(ctitle, {"size": 16, "bold": True, "color": ACCENT})], space_after=12, first=True)
    for name, desc in rows:
        para(tf, [(name, {"size": 14, "bold": True, "color": TEXT})], space_after=1, space_before=6)
        para(tf, [(desc, {"size": 11, "color": MUTED})], space_after=4)
    x += 3.83
tf = textbox(s, Inches(0.95), Inches(6.35), Inches(11.5), Inches(0.6))
para(tf, [("원칙  ", {"size": 12, "bold": True, "color": GREEN}),
          ("의존성 최소화(No Odin) · 외부 유료 에셋은 라이선스 보호 위해 별도 private repo 로 분리(gitignore)",
           {"size": 12, "color": MUTED})], first=True)
footer(s, pageno())

# =========================================================
# 4. 게임 소개
# =========================================================
s = slide(); header(s, "INTRODUCTION", "게임 소개 — “시간을 압축한 FM”")
tf = textbox(s, Inches(0.95), Inches(2.0), Inches(6.7), Inches(4.6))
para(tf, [("문제", {"size": 15, "bold": True, "color": AMBER})], space_after=3, first=True)
para(tf, [("기존 FM 은 깊지만 한 시즌에 수십 시간 → 직장인이 즐기기 부담", {"size": 14, "color": TEXT})], space_after=12, bullet="•")
para(tf, [("해법", {"size": 15, "bold": True, "color": GREEN})], space_after=3)
para(tf, [("경기 ‘시각화’를 과감히 포기하고 육성·이적·유스에 집중", {"size": 14, "color": TEXT})], space_after=2, bullet="•")
para(tf, [("한 시즌을 2~3시간에 완주", {"size": 14, "color": TEXT})], space_after=14, bullet="•")
para(tf, [("핵심 셀링 포인트", {"size": 15, "bold": True, "color": ACCENT})], space_after=4)
for t in ["3중 가챠 — 초기 스쿼드 / 유스 인스펙션 / 트레잇",
          "정보 비대칭 기반 의사결정 (제한된 정보 속 선택)",
          "리롤 ‘시스템화’ — 세이브로드 편법을 정식 메커닉으로",
          "시드 기반 리플레이 가치"]:
    para(tf, [(t, {"size": 13, "color": TEXT})], space_after=3, bullet="▹", level=0)
image_or_placeholder(s, "04_dashboard.png", 7.9, 2.0, 4.65, 4.5,
            "메인 대시보드 — 다음 경기 · 일정 · 인박스 허브")
footer(s, pageno())

# =========================================================
# 5. 구현 범위
# =========================================================
s = slide(); header(s, "SCOPE", "구현 범위 — 완성된 핵심 시스템")
sys_cards = [
    ("매치 엔진", "5-Zone Markov + xG + FM식 평점\n결정성(시드) · 텍스트 이벤트 · 슛맵", ACCENT),
    ("선수 · 성장", "49 stat · CA/PA · 신체 조건\n월별 성장 스냅샷 · 트레잇", GREEN),
    ("이적 시장", "스카우팅 · 협상 · 임대\n세부 stat / 주급 필터", ACCENT),
    ("유스", "정기 인스펙션 · 콜업\n멘토링", GREEN),
    ("전술 · 라인업", "포메이션 도식 드래그 · Role/Duty\n시너지 10종 · 포메이션 상성", ACCENT),
    ("운영", "시설 8종×10단계 · 재정\n시즌 사이클(시상/보드/리그)", GREEN),
    ("UX", "인박스 기반 알림 · GlobalNav\n옵션 · 사운드 · 통화", ACCENT),
    ("인프라", "세이브/로드 · 로컬라이즈(KO/EN)\n측정 하네스 · EditMode 테스트", GREEN),
]
x0, y0, w, h, gx, gy = 0.95, 1.95, 3.4, 0.98, 0.16, 0.15
for i, (t, d, ac) in enumerate(sys_cards):
    col = i % 2; row = i // 2
    l = x0 + col*(w+gx); tp = y0 + row*(h+gy)
    card(s, Inches(l), Inches(tp), Inches(w), Inches(h), accent=ac)
    tf = textbox(s, Inches(l+0.24), Inches(tp+0.1), Inches(w-0.4), Inches(h-0.18), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(t, {"size": 13.5, "bold": True, "color": TEXT})], space_after=3, first=True)
    para(tf, [(d, {"size": 10, "color": MUTED})], space_after=0, line=1.05)
image_or_placeholder(s, "05_squad_tactic.png", 8.2, 1.95, 4.25, 4.35,
            "전술 — 포메이션 피치 드래그 배정")
tf = textbox(s, Inches(0.95), Inches(6.5), Inches(11.5), Inches(0.5))
para(tf, [("범위 제외  ", {"size": 11.5, "bold": True, "color": AMBER}),
          ("챔피언스리그 · 국가대표 · 기자회견 · 코치진 관리 (1인 3주 스코프상 의도적 컷)",
           {"size": 11.5, "color": MUTED})], first=True)
footer(s, pageno())

# =========================================================
# 6. 개발 일정
# =========================================================
s = slide(); header(s, "TIMELINE", "개발 일정 — 3주, 3개의 마일스톤")
miles = [
    ("V0.1", "5/15–5/22", "프로토타입", "한 시즌을 끝까지 (UI 못생겨도 기능 우선)", ["코어 데이터 · 세이브/로드", "시간 진행 · 시즌 사이클", "스코어만 나오는 매치", "가챠 · 유스 · 단순 이적"], ACCENT),
    ("V0.5", "5/23–5/29", "플레이어빌리티", "‘재미있게’ 한 시즌이 굴러가게", ["텍스트 매치 이벤트 · 상세 스탯", "부상/카드 · 전술 프리셋", "사기/약속 · 스카우팅", "리그 3개 · 밸런싱"], GREEN),
    ("V1.0", "5/30–6/7", "폴리싱", "포트폴리오 빌드", ["GlobalNav · 인박스 UX 패러다임", "xG 매치 재설계 · 재정 밸런싱", "옵션/사운드/통화 · 크레스트", "Unity MCP 도입 · 측정 튜닝"], AMBER),
]
x = 0.95
for tag, period, title, goal, rows, ac in miles:
    card(s, Inches(x), Inches(1.95), Inches(3.78), Inches(4.05), accent=ac)
    tf = textbox(s, Inches(x+0.3), Inches(2.18), Inches(3.2), Inches(3.7))
    para(tf, [(tag, {"size": 26, "bold": True, "color": ac})], space_after=0, first=True)
    para(tf, [(period + "  ·  " + title, {"size": 12, "bold": True, "color": TEXT})], space_after=8)
    para(tf, [(goal, {"size": 11.5, "italic": True, "color": MUTED})], space_after=10, line=1.15)
    for r in rows:
        para(tf, [(r, {"size": 12, "color": TEXT})], space_after=4, bullet="·", level=0)
    x += 3.95
tf = textbox(s, Inches(0.95), Inches(6.25), Inches(11.5), Inches(0.7))
para(tf, [("방법론  ", {"size": 12, "bold": True, "color": ACCENT}),
          ("이슈→브랜치→PR→머지 · Sub-PR 패턴(명세→구현→에셋 chore) · ", {"size": 12, "color": MUTED}),
          ("‘명세 우선’ — 설계 결정 77건을 design-decisions.md 에 기록", {"size": 12, "bold": True, "color": TEXT})],
     first=True, line=1.2)
footer(s, pageno())

# =========================================================
# 7. 아키텍처
# =========================================================
s = slide(); header(s, "ARCHITECTURE", "아키텍처 — 단방향 4-Layer")
layers = [
    ("Presentation", "UI 컨트롤러 · Unity Scene", RGBColor(0x35,0x44,0x6B)),
    ("Application", "Stateless Systems — MatchSimulator · TransferSystem · GrowthSystem …", RGBColor(0x2C,0x3A,0x5E)),
    ("Core", "GameManager(진입점) · EventBus · GameTime — 인프라", RGBColor(0x24,0x31,0x50)),
    ("Domain", "Player · Club · Match … 게임 규칙·데이터  (외부 의존 0)", RGBColor(0x1C,0x27,0x42)),
    ("Data", "Save/Load · ScriptableObject", RGBColor(0x16,0x1F,0x36)),
]
y = 1.95
for i,(name, desc, col) in enumerate(layers):
    inset = Inches(0.0)
    card(s, Inches(2.2), Inches(y), Inches(7.0), Inches(0.78), fill=col, accent=ACCENT)
    tf = textbox(s, Inches(2.5), Inches(y), Inches(6.6), Inches(0.78), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(name, {"size": 15, "bold": True, "color": TEXT}),
              ("    " + desc, {"size": 11, "color": MUTED})], first=True, space_after=0)
    if i < len(layers)-1:
        ar = rect(s, Inches(5.6), Inches(y+0.78), Inches(0.2), Inches(0.1), fill=ACCENT, shape=MSO_SHAPE.DOWN_ARROW)
    y += 0.93
# 원칙 사이드
tf = textbox(s, Inches(9.5), Inches(2.05), Inches(3.1), Inches(4.3))
para(tf, [("설계 원칙", {"size": 15, "bold": True, "color": ACCENT})], space_after=8, first=True)
for t in ["의존 방향 단방향\n— 역참조·순환 금지", "Stateless 시스템\n— GameState 입출력만",
          "ID 기반 참조\n— 직접 참조 금지 → 직렬화 안전", "매직넘버 0\n— GameBalanceSO 외부화"]:
    para(tf, [(t, {"size": 12, "color": TEXT})], space_after=9, bullet="•", line=1.12)
footer(s, pageno())

# =========================================================
# 8. 설계가 바뀐 순간들
# =========================================================
s = slide(); header(s, "PROCESS · PIVOTS", "설계가 바뀐 순간들")
pivots = [
    ("강제 씬 전환  →  인박스 패러다임",
     "초기엔 이벤트마다 씬을 강제 이동 → ‘자유도 침해’ 피드백.\nInboxRouter 가 이벤트를 알림으로 흡수, 유저가 직접 선택하게 전환."),
    ("블런트 튜닝  →  xG 찬스-퀄리티 레이어",
     "‘4-파라미터 튜닝’으로 골 수를 맞추려다 한계.\n더 FM다운 방식을 찾아 실제 매치엔진 원리(xG)를 도입."),
    ("가격 상향  →  현금흐름 정상화",
     "‘시설비만 올리자’는 표면 처방을 버리고 근본 원인 진단.\n상시 유출(주급) 부재가 핵심 → 월 주급 차감 + 수입 비율 정상화."),
]
y = 2.1
for t, d in pivots:
    card(s, Inches(0.95), Inches(y), Inches(11.4), Inches(1.3), accent=ACCENT)
    tf = textbox(s, Inches(1.3), Inches(y+0.16), Inches(10.8), Inches(1.0), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(t, {"size": 16, "bold": True, "color": TEXT})], space_after=3, first=True)
    para(tf, [(d, {"size": 12.5, "color": MUTED})], space_after=0, line=1.15)
    y += 1.5
footer(s, pageno())

# =========================================================
# 9. 난관 ① 게임플레이 밸런싱
# =========================================================
s = slide(); header(s, "CHALLENGES", "난관과 극복 — 게임플레이 밸런싱")
rows = [
    ("과득점", "박스 도달 = 동급 찬스 → 경기당 평균 5~6골",
     "xG 모델로 ‘찬스 품질’을 분리. E[goals]=ΣxG 로 수학적 산정 → 측정 평균 2.76골 (목표 2.7±0.3)", GREEN),
    ("재정 무한 증식", "주급 유출 0 + 수입만 누적 → 자금이 영구 증가",
     "월 주급 차감 신설 + EPL 임금/매출 63% 앵커로 수입 재스케일 → 중위 구단 본전±", GREEN),
    ("비현실 평점", "포지션 무가중 → 수비/미드가 항상 저평점",
     "각 포지션 액션을 충분히 가치화 + xG 보정(빅찬스 미스 급락)으로 FM 정합", GREEN),
]
y = 1.98
for t, prob, sol, ac in rows:
    card(s, Inches(0.95), Inches(y), Inches(11.4), Inches(1.18))
    rect(s, Inches(0.95), Inches(y), Inches(2.3), Inches(1.18), fill=BG_CODE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf0 = textbox(s, Inches(1.05), Inches(y), Inches(2.1), Inches(1.18), anchor=MSO_ANCHOR.MIDDLE)
    para(tf0, [(t, {"size": 15, "bold": True, "color": AMBER})], align=PP_ALIGN.CENTER, first=True)
    tf = textbox(s, Inches(3.45), Inches(y+0.16), Inches(8.7), Inches(0.95), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [("문제  ", {"size": 11.5, "bold": True, "color": MUTED}), (prob, {"size": 12.5, "color": TEXT})], space_after=4, first=True)
    para(tf, [("해결  ", {"size": 11.5, "bold": True, "color": GREEN}), (sol, {"size": 12.5, "color": TEXT})], space_after=0, line=1.08)
    y += 1.32
tf = textbox(s, Inches(0.95), Inches(6.55), Inches(11.5), Inches(0.5))
para(tf, [("공통 무기  ", {"size": 12, "bold": True, "color": ACCENT}),
          ("측정 하네스(200~380 매치 / 20구단 시즌 순현금흐름) + 사용자 Test Runner 반복 튜닝", {"size": 12, "color": MUTED})], first=True)
footer(s, pageno())

# =========================================================
# 10. 느낀 점
# =========================================================
s = slide(); header(s, "REFLECTION", "느낀 점 — 3주가 남긴 것")
notes = [
    ("‘명세 우선’의 힘", "즉흥 결정 대신 설계 결정 77건을 문서화 → 일관성과 ‘되돌림’ 방지. 혼자여도 미래의 나를 위한 계약서."),
    ("밸런싱은 직감이 아닌 수학", "xG = ΣE[goals] 처럼 목표를 수식으로 환원하니 시행착오가 1~2회로 수렴."),
    ("‘버리는 결정’이 완성도를 만든다", "경기 시각화 포기 같은 과감한 컷이 1인 3주 스코프를 끝까지 가게 했다."),
    ("AI 협업의 분업", "설계·판단은 사람, 반복 구현·리팩터·문서·와이어링은 AI. 혼자선 못 했을 범위를 3주에 완주."),
]
y = 2.15
for t, d in notes:
    card(s, Inches(0.95), Inches(y), Inches(11.4), Inches(0.95), accent=GREEN)
    tf = textbox(s, Inches(1.3), Inches(y), Inches(10.9), Inches(0.95), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(t, {"size": 15, "bold": True, "color": GREEN}), ("    " + d, {"size": 12.5, "color": TEXT})], first=True, line=1.12)
    y += 1.08
footer(s, pageno())

# =========================================================
# 12. 대표 구현 — 매치 개요
# =========================================================
s = slide(); header(s, "DEEP DIVE · 매치 시뮬레이션 1/3", "대표 구현 ① — 매치 엔진 개요")
tf = textbox(s, Inches(0.95), Inches(2.0), Inches(6.7), Inches(4.5))
para(tf, [("경기 시각화를 포기한 대신, 결과를 만드는 ‘엔진’에 집중", {"size": 13.5, "italic": True, "color": MUTED})], space_after=12, first=True)
for lead, t in [("5-Zone Markov", "DefThird→Mid→AttThird→Box→Shot 의 확률 전이로 한 경기를 이벤트 시퀀스로 전개"),
                ("결정성(시드)", "매치 시작 시 시드 고정 → 같은 시드 = 같은 경기. ‘리롤 시스템’의 기술적 근간"),
                ("동일 엔진", "유저 경기와 백그라운드 리그 경기를 같은 엔진으로 → 리그 전체 일관성"),
                ("풍부한 출력", "스코어 + 텍스트 이벤트 + 슛맵/히트맵 + 선수별 평점·스탯")]:
    para(tf, [(lead, {"size": 14, "bold": True, "color": ACCENT, "bcolor": ACCENT})], space_after=2, space_before=4, bullet="●")
    para(tf, [(t, {"size": 12.5, "color": TEXT})], space_after=6, line=1.12)
image_or_placeholder(s, "11_match_text.png", 7.9, 2.0, 4.6, 4.5,
            "매치 텍스트 씬 — 실시간 이벤트 · 가속(1/2/3/4)")
footer(s, pageno())

# =========================================================
# 13. 매치 — xG 레이어
# =========================================================
s = slide(); header(s, "DEEP DIVE · 매치 시뮬레이션 2/3", "xG 찬스-퀄리티 레이어")
tf = textbox(s, Inches(0.95), Inches(1.95), Inches(6.6), Inches(4.6))
para(tf, [("핵심 아이디어 — “슛의 가치 = 그 찬스가 어떻게 만들어졌나”", {"size": 14, "bold": True, "color": ACCENT})], space_after=10, first=True)
for lead, t in [("chanceType 별 baseXG", "ClearChance · OpenPlay · Header · LongShot · FreeKick · Penalty"),
                ("기록되는 xG", "= 찬스 품질(situation). 슈터 실력과 무관 (표준 xG 정의)"),
                ("실제 골", "= xG × finishMod(슈터 결정력) × gkMod(GK 보정)"),
                ("밸런싱이 ‘수학’", "E[goals] ≈ ΣxG → 팀당 ΣxG≈1.35 → 2.7/경기 직접 산정")]:
    para(tf, [(lead, {"size": 13.5, "bold": True, "color": TEXT, "bcolor": ACCENT})], space_after=1, space_before=3, bullet="▸")
    para(tf, [(t, {"size": 12, "color": MUTED})], space_after=5, line=1.1)
# 슛맵 (xG 시각 증거)
image_or_placeholder(s, "12_shotmap.png", 7.85, 1.95, 4.55, 2.35,
            "매치 결과 — 슛맵 / xG 시각화")
# 측정 결과
card(s, Inches(7.85), Inches(4.5), Inches(4.55), Inches(1.9), accent=GREEN)
tfm = textbox(s, Inches(8.15), Inches(4.7), Inches(4.0), Inches(1.6))
para(tfm, [("측정 결과 (200 매치)", {"size": 13, "bold": True, "color": GREEN})], space_after=6, first=True)
para(tfm, [("평균 2.76골", {"size": 13, "bold": True, "color": TEXT}), ("   ·   std 1.62   ·   PK 0.26/경기", {"size": 12, "color": MUTED})], space_after=3, first=False)
para(tfm, [("목표 2.7±0.3 정합 ✓", {"size": 12, "color": GREEN})], space_after=0)
footer(s, pageno())

# =========================================================
# 14. 매치 — 평점 재설계
# =========================================================
s = slide(); header(s, "DEEP DIVE · 매치 시뮬레이션 3/3", "평점 재설계 — FM 정합")
tf = textbox(s, Inches(0.95), Inches(2.0), Inches(6.6), Inches(4.4))
para(tf, [("문제", {"size": 14, "bold": True, "color": AMBER})], space_after=3, first=True)
para(tf, [("포지션 무가중 → 수비/미드 저평점 · 실점 책임을 GK 가 독점 · 부진 감점 부족", {"size": 12.5, "color": TEXT})], space_after=12, bullet="•", line=1.15)
para(tf, [("원리 — “포지션이 평점을 만든다”", {"size": 14, "bold": True, "color": ACCENT})], space_after=4)
for t in ["인위적 라인 곱셈 대신, 각 포지션 액션을 충분히 가치화",
          "패스 성공률 티어 보너스 · 수비 액션/클리어런스 가치 ↑",
          "DF 무실점/실점 공유 · GK 독점 해소",
          "xG 보정 — 빅찬스 미스 급락 / clinical finish 가산",
          "9개 평점 파라미터를 GameBalanceSO 로 외부화"]:
    para(tf, [(t, {"size": 12.5, "color": TEXT})], space_after=4, bullet="▹")
# 예시(컴팩트) + 평점 화면
card(s, Inches(7.85), Inches(2.0), Inches(4.55), Inches(1.78))
tfe = textbox(s, Inches(8.12), Inches(2.16), Inches(4.05), Inches(1.5), anchor=MSO_ANCHOR.MIDDLE)
para(tfe, [("xG 보정 직관", {"size": 13, "bold": True, "color": ACCENT})], space_after=6, first=True)
para(tfe, [("0.40 xG 빅찬스 미스   →   ", {"size": 12, "color": TEXT}), ("−0.74", {"size": 14, "bold": True, "color": AMBER})], space_after=3)
para(tfe, [("0.05 xG 침착한 골   →   ", {"size": 12, "color": TEXT}), ("+1.37", {"size": 14, "bold": True, "color": GREEN})], space_after=0)
image_or_placeholder(s, "13_ratings.png", 7.85, 3.95, 4.55, 2.45,
            "매치 결과 — 포지션별 선수 평점")
footer(s, pageno())

# =========================================================
# 15. 정보 비대칭 & 스카우팅
# =========================================================
s = slide(); header(s, "DEEP DIVE · 기획-기술 결합", "대표 구현 ② — 정보 비대칭 & 스카우팅")
tf = textbox(s, Inches(0.95), Inches(2.0), Inches(6.6), Inches(4.4))
para(tf, [("게임의 핵심 재미 = 제한된 정보 속 의사결정. ‘정보’ 자체를 시스템으로 설계", {"size": 13, "italic": True, "color": MUTED})], space_after=12, first=True)
for lead, t in [("5단계 티어 표시", "정확한 수치는 숨기고, 명성 대비 상대평가로 등급만 노출"),
                ("Trait 가시성 3-tier", "Concealed(영구 비노출) / Public(항상) / ScoutGated(정찰 비례)"),
                ("스카우팅 게이팅", "scoutLevel(0~100) 비례로 공개 trait 수 결정 — ‘검색으로 스카우팅 무력화’ 차단"),
                ("정보 누설 방지", "미정찰 선수는 보유 trait 수조차 숨김 → ‘추가 정찰 필요’만 표시")]:
    para(tf, [(lead, {"size": 13.5, "bold": True, "color": TEXT, "bcolor": ACCENT})], space_after=2, space_before=4, bullet="●")
    para(tf, [(t, {"size": 12, "color": MUTED})], space_after=4, line=1.1)
# 선수 프로필 화면 + 3-tier 레전드
image_or_placeholder(s, "14_player_profile.png", 7.85, 1.98, 4.55, 2.3,
            "선수 프로필 — 스탯 색상 등급 + 미정찰 trait 가림")
card(s, Inches(7.85), Inches(4.55), Inches(4.55), Inches(1.85), accent=ACCENT)
tfl = textbox(s, Inches(8.12), Inches(4.72), Inches(4.05), Inches(1.6))
para(tfl, [("Trait 가시성 3-tier", {"size": 13, "bold": True, "color": ACCENT})], space_after=7, first=True)
para(tfl, [("Concealed", {"size": 12, "bold": True, "color": AMBER}), ("  비노출 (늦깎이형·조숙형)", {"size": 11, "color": MUTED})], space_after=4)
para(tfl, [("Public", {"size": 12, "bold": True, "color": GREEN}), ("  항상 (골결정력·만능형 …)", {"size": 11, "color": MUTED})], space_after=4)
para(tfl, [("ScoutGated", {"size": 12, "bold": True, "color": ACCENT}), ("  정찰 비례 (부상취약 …)", {"size": 11, "color": MUTED})], space_after=0)
footer(s, pageno())

# =========================================================
# 16. AI 협업 워크플로우
# =========================================================
s = slide(); header(s, "DEEP DIVE · 프로세스", "대표 구현 ③ — AI 협업 워크플로우")
tf = textbox(s, Inches(0.95), Inches(1.95), Inches(11.5), Inches(0.6))
para(tf, [("1인 개발이지만 ‘혼자’가 아니었다 — 역할을 나눈 도구 오케스트레이션", {"size": 13.5, "italic": True, "color": MUTED})], first=True)
roles = [
    ("채팅 (Claude)", "설계 · 디자인 결정\n알고리즘 명세", ACCENT),
    ("Claude Code + Unity MCP", "코드 · git · PR\n씬/prefab/인스펙터 직접 와이어링", GREEN),
    ("GitHub Issues / Projects", "작업 추적\n이슈 → 브랜치 → PR → 보드", ACCENT),
]
x = 0.95
for t, d, ac in roles:
    card(s, Inches(x), Inches(2.7), Inches(3.66), Inches(1.85), accent=ac)
    tfr = textbox(s, Inches(x+0.28), Inches(2.92), Inches(3.1), Inches(1.5))
    para(tfr, [(t, {"size": 14.5, "bold": True, "color": TEXT})], space_after=6, first=True, line=1.05)
    para(tfr, [(d, {"size": 12, "color": MUTED})], space_after=0, line=1.15)
    x += 3.87
tf = textbox(s, Inches(0.95), Inches(4.85), Inches(7.0), Inches(1.7))
para(tf, [("기록 기반 협업", {"size": 14, "bold": True, "color": ACCENT})], space_after=6, first=True)
para(tf, [("CLAUDE.md · design-decisions.md · 영속 메모리로 세션 간 컨텍스트를 유지 — ‘즉흥 결정 금지’를 시스템화",
           {"size": 12.5, "color": TEXT})], space_after=12, bullet="•", line=1.18)
para(tf, [("결과  ", {"size": 13, "bold": True, "color": GREEN}),
          ("1인 · 3주에 설계 결정 77건 / 50+ PR / 4-Layer 풀시스템", {"size": 13, "bold": True, "color": TEXT})], space_after=0)
image_or_placeholder(s, "15_github_board.png", 8.25, 4.78, 4.15, 1.95,
            "GitHub Projects — 이슈→PR→머지 추적")
footer(s, pageno())

# =========================================================
# 17. 결과 & 회고
# =========================================================
s = slide(); header(s, "WRAP-UP", "결과 & 회고")
cols = [
    ("완성한 것", GREEN, ["한 시즌 풀 플레이 루프\n(가챠→육성→이적→유스→전술→매치→결산)",
                       "FM다운 매치 엔진 (xG + 평점)",
                       "정보 비대칭 메타 (스카우팅 게이팅)",
                       "V0.1→V1.0 3주 마일스톤 달성"]),
    ("핵심 자산", ACCENT, ["기록 기반 개발 프로세스\n(명세 우선 · 측정 하네스)",
                       "재사용 가능한 측정/튜닝 루프",
                       "AI 협업 워크플로우 노하우",
                       "4-Layer 확장 가능 구조"]),
    ("남은 것 (V1.x)", AMBER, ["컵 대회 · 세이브 슬롯명",
                           "남은 플레이테스트 핫픽스(Stage R)",
                           "xG 좌표 정밀화 · Role별 평점",
                           "사운드/로컬라이즈 폴리싱"]),
]
x = 0.95
for t, ac, rows in cols:
    card(s, Inches(x), Inches(2.0), Inches(3.78), Inches(4.3), accent=ac)
    tf = textbox(s, Inches(x+0.3), Inches(2.25), Inches(3.2), Inches(3.9))
    para(tf, [(t, {"size": 16, "bold": True, "color": ac})], space_after=12, first=True)
    for r in rows:
        para(tf, [(r, {"size": 12, "color": TEXT})], space_after=8, bullet="•", line=1.12)
    x += 3.95
footer(s, pageno())

# =========================================================
# 18. Q&A
# =========================================================
s = slide()
rect(s, Inches(8.8), Inches(-1.5), Inches(6), Inches(6), fill=BG_CARD, shape=MSO_SHAPE.OVAL)
rect(s, Inches(1.1), Inches(2.5), Inches(0.16), Inches(2.3), fill=ACCENT)
tf = textbox(s, Inches(1.5), Inches(2.6), Inches(9), Inches(2.6))
para(tf, [("Q & A", {"size": 80, "bold": True, "color": TEXT})], space_after=6, first=True)
para(tf, [("들어주셔서 감사합니다.", {"size": 22, "color": MUTED})], space_after=0)
tf2 = textbox(s, Inches(1.5), Inches(5.6), Inches(10), Inches(0.6))
para(tf2, [("FM-Lite", {"size": 14, "bold": True, "color": ACCENT}),
           ("   ·   JihooKim   ·   jjihooo49@gmail.com", {"size": 14, "color": MUTED})], first=True)

import os
out = os.path.join(os.path.dirname(__file__), "FM-Lite_발표자료.pptx")
prs.save(out)
print("saved:", out, "| slides:", len(prs.slides._sldIdLst))
